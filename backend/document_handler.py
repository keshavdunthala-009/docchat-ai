import os
import re
from typing import List, Dict
from sentence_transformers import SentenceTransformer
import chromadb


class DocumentProcessor:

    def __init__(self, session_id: str = "default"):
        self.model = SentenceTransformer('all-MiniLM-L6-v2')
        self.client = chromadb.PersistentClient(path="/tmp/chroma_data")
        self.collection = self.client.get_or_create_collection(
            name=f"documents_{session_id}",
            metadata={"hnsw:space": "cosine"}
        )

    # ============================================================
    # MAIN EXTRACTION METHOD
    # ============================================================
    def extract_text_from_pdf(self, file_path: str) -> str:
        """Main method - detects file type and extracts text"""
        
        extension = file_path.lower().split('.')[-1]
        print(f"File type detected: {extension}")
        
        if extension == 'pdf':
            return self.extract_from_pdf(file_path)
        elif extension in ['xlsx', 'xls']:
            return self.extract_from_excel(file_path)
        elif extension == 'csv':
            return self.extract_from_csv(file_path)
        elif extension == 'txt':
            return self.extract_from_txt(file_path)
        elif extension == 'docx':
            return self.extract_from_word(file_path)
        elif extension in ['pptx', 'ppt']:
            return self.extract_from_powerpoint(file_path)
        elif extension in ['png', 'jpg', 'jpeg', 'tiff', 'bmp']:
            return self.extract_from_image(file_path)
        else:
            return self.extract_from_txt(file_path)

    # ============================================================
    # PDF EXTRACTION
    # ============================================================
    def extract_from_pdf(self, pdf_path: str) -> str:
        """Extract from PDF - handles both text and scanned PDFs"""
        
        # First try pdfplumber (best for text + tables)
        try:
            import pdfplumber
            text = ""
            seen_content = set()

            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    # Extract regular text
                    page_text = page.extract_text()
                    if page_text and page_text not in seen_content:
                        text += page_text + "\n"
                        seen_content.add(page_text)

                    # Extract tables
                    tables = page.extract_tables()
                    for table in tables:
                        if table:
                            table_text = "\n[TABLE]\n"
                            for row in table:
                                if row:
                                    clean_row = [
                                        str(cell).strip() if cell else ""
                                        for cell in row
                                    ]
                                    row_text = " | ".join(clean_row)
                                    if row_text.strip():
                                        table_text += row_text + "\n"
                            table_text += "[/TABLE]\n"
                            if table_text not in seen_content:
                                text += table_text
                                seen_content.add(table_text)

            if text.strip():
                print(f"pdfplumber extracted: {len(text)} chars")
                return text

        except Exception as e:
            print(f"pdfplumber failed: {e}")

        # Second try PyPDF2
        try:
            import PyPDF2
            text = ""
            with open(pdf_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    extracted = page.extract_text()
                    if extracted:
                        text += extracted + "\n"

            if text.strip():
                print(f"PyPDF2 extracted: {len(text)} chars")
                return text

        except Exception as e:
            print(f"PyPDF2 failed: {e}")

        # Last resort - OCR for scanned PDFs
        print("Trying OCR for scanned PDF...")
        return self.extract_from_scanned_pdf(pdf_path)

    # ============================================================
    # SCANNED PDF / IMAGE OCR
    # ============================================================
    def extract_from_scanned_pdf(self, pdf_path: str) -> str:
        """Extract text from scanned PDF using OCR"""
        try:
            from pdf2image import convert_from_path
            import pytesseract

            text = ""
            print("Converting PDF pages to images...")
            
            images = convert_from_path(pdf_path, dpi=300)
            print(f"Converted {len(images)} pages")

            for i, image in enumerate(images):
                print(f"OCR processing page {i+1}...")
                page_text = pytesseract.image_to_string(
                    image,
                    lang='eng',
                    config='--psm 6'
                )
                if page_text.strip():
                    text += f"\n[PAGE {i+1}]\n{page_text}\n"

            if text.strip():
                print(f"OCR extracted: {len(text)} chars")
                return text
            else:
                return "Could not extract text from this scanned PDF"

        except Exception as e:
            print(f"OCR failed: {e}")
            return "OCR extraction failed. Please ensure the PDF is readable."

    def extract_from_image(self, image_path: str) -> str:
        """Extract text from image using OCR"""
        try:
            import pytesseract
            from PIL import Image

            print(f"OCR on image: {image_path}")
            image = Image.open(image_path)
            text = pytesseract.image_to_string(
                image,
                lang='eng',
                config='--psm 6'
            )

            if text.strip():
                print(f"Image OCR extracted: {len(text)} chars")
                return text
            else:
                return "Could not extract text from image"

        except Exception as e:
            print(f"Image OCR failed: {e}")
            return f"Image extraction failed: {str(e)}"

    # ============================================================
    # EXCEL EXTRACTION
    # ============================================================
    def extract_from_excel(self, excel_path: str) -> str:
        """Extract from Excel - reads ALL sheets"""
        try:
            import openpyxl
            text = ""

            workbook = openpyxl.load_workbook(excel_path)
            print(f"Excel sheets: {workbook.sheetnames}")

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"\n[SHEET: {sheet_name}]\n"

                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        clean_row = [
                            str(cell).strip() if cell is not None else ""
                            for cell in row
                        ]
                        text += " | ".join(clean_row) + "\n"

                text += f"[/SHEET]\n"

            print(f"Excel extracted: {len(text)} chars")
            return text

        except Exception as e:
            print(f"Excel extraction failed: {e}")
            return f"Excel extraction failed: {str(e)}"

    # ============================================================
    # CSV EXTRACTION
    # ============================================================
    def extract_from_csv(self, csv_path: str) -> str:
        """Extract from CSV"""
        try:
            import csv
            text = "[CSV DATA]\n"

            with open(csv_path, 'r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    if any(cell.strip() for cell in row):
                        text += " | ".join(row) + "\n"

            text += "[/CSV DATA]\n"
            print(f"CSV extracted: {len(text)} chars")
            return text

        except Exception as e:
            print(f"CSV extraction failed: {e}")
            return f"CSV extraction failed: {str(e)}"

    # ============================================================
    # WORD EXTRACTION
    # ============================================================
    def extract_from_word(self, docx_path: str) -> str:
        """Extract from Word document"""
        try:
            from docx import Document
            text = ""

            doc = Document(docx_path)

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"

            # Extract tables
            for table in doc.tables:
                text += "\n[TABLE]\n"
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    text += " | ".join(row_data) + "\n"
                text += "[/TABLE]\n"

            print(f"Word extracted: {len(text)} chars")
            return text

        except Exception as e:
            print(f"Word extraction failed: {e}")
            return f"Word extraction failed: {str(e)}"

    # ============================================================
    # POWERPOINT EXTRACTION
    # ============================================================
    def extract_from_powerpoint(self, pptx_path: str) -> str:
        """Extract from PowerPoint"""
        try:
            from pptx import Presentation
            text = ""

            prs = Presentation(pptx_path)
            print(f"PowerPoint slides: {len(prs.slides)}")

            for i, slide in enumerate(prs.slides):
                text += f"\n[SLIDE {i+1}]\n"

                for shape in slide.shapes:
                    # Extract text from shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

                    # Extract tables from slides
                    if shape.has_table:
                        text += "\n[TABLE]\n"
                        for row in shape.table.rows:
                            row_data = [
                                cell.text.strip()
                                for cell in row.cells
                            ]
                            text += " | ".join(row_data) + "\n"
                        text += "[/TABLE]\n"

                text += f"[/SLIDE]\n"

            print(f"PowerPoint extracted: {len(text)} chars")
            return text

        except Exception as e:
            print(f"PowerPoint extraction failed: {e}")
            return f"PowerPoint extraction failed: {str(e)}"

    # ============================================================
    # TXT EXTRACTION
    # ============================================================
    def extract_from_txt(self, txt_path: str) -> str:
        """Extract from text file"""
        try:
            with open(txt_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        except Exception as e:
            print(f"TXT extraction failed: {e}")
            return ""

    # ============================================================
    # CHUNKING
    # ============================================================
    def chunk_text(self, text: str, chunk_size: int = 500, overlap: int = 100) -> List[Dict]:
        chunks = []
        chunk_id = 0

        sections = re.split(r'\n(?=[A-Z][A-Z\s]+\n)', text)

        for section in sections:
            section = section.strip()
            if len(section) > 50:
                if len(section) > chunk_size:
                    start = 0
                    while start < len(section):
                        end = start + chunk_size
                        chunk = section[start:end].strip()
                        if len(chunk) > 50:
                            chunks.append({
                                "id": f"chunk_{chunk_id}",
                                "text": chunk,
                                "chunk_index": chunk_id
                            })
                            chunk_id += 1
                        start = end - overlap
                else:
                    chunks.append({
                        "id": f"chunk_{chunk_id}",
                        "text": section,
                        "chunk_index": chunk_id
                    })
                    chunk_id += 1

        if len(chunks) == 0:
            start = 0
            while start < len(text):
                end = start + chunk_size
                chunk = text[start:end].strip()
                if len(chunk) > 50:
                    chunks.append({
                        "id": f"chunk_{chunk_id}",
                        "text": chunk,
                        "chunk_index": chunk_id
                    })
                    chunk_id += 1
                start = end - overlap

        return chunks

    # ============================================================
    # STORE DOCUMENT
    # ============================================================
    def store_document(self, pdf_path: str, document_name: str) -> dict:
        print(f"Processing: {document_name}")
        text = self.extract_text_from_pdf(pdf_path)
        print(f"Extracted {len(text)} characters")
        chunks = self.chunk_text(text)
        print(f"Created {len(chunks)} chunks")
        chunk_texts = [chunk["text"] for chunk in chunks]

        embeddings = self.model.encode(
            chunk_texts,
            batch_size=32,
            convert_to_numpy=True
        )
        print(f"Embedded {len(embeddings)} chunks")

        ids = [f"{document_name}_chunk_{i}" for i in range(len(chunks))]
        metadatas = [{"document": document_name, "chunk_index": i} for i in range(len(chunks))]

        self.collection.add(
            ids=ids,
            embeddings=embeddings.tolist(),
            documents=chunk_texts,
            metadatas=metadatas
        )
        print(f"Stored all {len(chunks)} chunks!")
        return {"chunks": len(chunks), "document": document_name}

    # ============================================================
    # SEARCH DOCUMENTS
    # ============================================================
    def search_documents(self, query: str, top_k: int = 7) -> List[Dict]:
        query_embedding = self.model.encode(query)
        results = self.collection.query(
            query_embeddings=[query_embedding.tolist()],
            n_results=top_k
        )
        search_results = []
        for i, doc in enumerate(results['documents'][0]):
            search_results.append({
                "text": doc,
                "metadata": results['metadatas'][0][i],
                "rank": i + 1
            })
        return search_results